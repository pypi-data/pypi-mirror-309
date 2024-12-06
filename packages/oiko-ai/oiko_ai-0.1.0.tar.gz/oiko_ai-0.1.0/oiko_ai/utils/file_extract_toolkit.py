from io import BytesIO

import docx
import fitz
import pandas as pd
from fastapi import UploadFile
from pptx import Presentation


def extract_text(
    file: UploadFile, allowed_extensions: list[str] = ['all']
) -> str:
    """
    Extrai texto de um arquivo recebido via UploadFile do FastAPI, com base nas extensões permitidas.

    Tipos de arquivos suportados:

    - **PDF**: Arquivos `.pdf`.
    - **Excel**: Arquivos `.xls` e `.xlsx`.
    - **Texto**: Arquivos `.txt`.
    - **PowerPoint**: Arquivos `.ppt` e `.pptx`.
    - **Word**: Arquivos `.doc` e `.docx`.

    Args:
        file (UploadFile): O arquivo enviado para upload.
        allowed_extensions (list[str]): Lista de extensões permitidas (ex.: `["pdf", "txt"]`) ou `["all"]` para permitir todos os tipos suportados.

    Returns:
        str: O texto extraído do arquivo.

    Raises:
        ValueError: Se o tipo de arquivo não for suportado ou não estiver na lista permitida.
    """
    # Tipos suportados
    supported_extensions = [
        'pdf',
        'xls',
        'xlsx',
        'txt',
        'ppt',
        'pptx',
        'doc',
        'docx',
    ]

    # Obtém a extensão do arquivo
    extension = file.filename.rsplit('.', 1)[-1].lower()

    # Validações
    if allowed_extensions != ['all'] and extension not in allowed_extensions:
        raise ValueError(f"Tipo de arquivo '{extension}' não permitido.")
    if extension not in supported_extensions:
        raise ValueError(f"Tipo de arquivo '{extension}' não suportado.")

    # Lê o conteúdo do arquivo como BytesIO
    file_stream = BytesIO(file.file.read())

    # Processa o arquivo com base em sua extensão
    if extension == 'pdf':
        return extract_text_from_pdf(file_stream)
    elif extension in ['xls', 'xlsx']:
        return extract_text_from_excel(file_stream)
    elif extension == 'txt':
        return extract_text_from_txt(file_stream)
    elif extension in ['ppt', 'pptx']:
        return extract_text_from_powerpoint(file_stream)
    elif extension in ['doc', 'docx']:
        return extract_text_from_doc(file_stream)
    else:
        raise ValueError(f"Tipo de arquivo '{extension}' não suportado.")


def extract_text_from_pdf(file_stream: BytesIO) -> str:
    text = ''
    pdf_document = fitz.open(stream=file_stream, filetype='pdf')
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text += page.get_text()
    return text


def extract_text_from_excel(file_stream: BytesIO) -> str:
    df = pd.read_excel(file_stream)
    text = df.head(100).to_string(index=False)
    return text


def extract_text_from_txt(file_stream: BytesIO) -> str:
    return file_stream.getvalue().decode('utf-8')


def extract_text_from_powerpoint(file_stream: BytesIO) -> str:
    text = ''
    presentation = Presentation(file_stream)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + '\n'
    return text


def extract_text_from_doc(file_stream: BytesIO) -> str:
    text = ''
    doc = docx.Document(file_stream)
    for para in doc.paragraphs:
        text += para.text + '\n'
    return text
